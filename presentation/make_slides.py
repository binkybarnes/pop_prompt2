from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
TEAL   = RGBColor(0x00, 0x87, 0x8A)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # truly blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def txbox(slide, text, x, y, w, h, size=18, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.15), fill=NAVY)
    txbox(slide, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.65),
          size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12), Inches(0.38),
              size=14, color=GOLD, align=PP_ALIGN.LEFT)


def footer(slide, text="Hack the Coast 2026 · Prince of Peace Enterprises"):
    rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=NAVY)
    txbox(slide, text, Inches(0.3), H - Inches(0.33), Inches(12), Inches(0.3),
          size=10, color=WHITE, align=PP_ALIGN.LEFT)


def bullet_box(slide, items, x, y, w, h, size=16, indent=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        prefix = "    • " if indent else "• "
        run.text = prefix + item
        run.font.size = Pt(size)
        run.font.color.rgb = DGRAY


def simple_table(slide, headers, rows, x, y, col_widths, row_h=Inches(0.42),
                 header_fill=NAVY, alt_fill=LGRAY):
    n_cols = len(headers)
    total_w = sum(col_widths)

    # header row
    cx = x
    for i, h_text in enumerate(headers):
        rect(slide, cx, y, col_widths[i], row_h, fill=header_fill)
        txbox(slide, h_text, cx + Inches(0.08), y + Inches(0.06),
              col_widths[i] - Inches(0.1), row_h - Inches(0.08),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        cx += col_widths[i]

    # data rows
    for r_idx, row in enumerate(rows):
        ry = y + row_h * (r_idx + 1)
        fill = alt_fill if r_idx % 2 == 0 else WHITE
        cx = x
        for c_idx, cell in enumerate(row):
            rect(slide, cx, ry, col_widths[c_idx], row_h, fill=fill)
            cell_color = DGRAY
            if isinstance(cell, tuple):
                cell, cell_color = cell
            txbox(slide, str(cell), cx + Inches(0.08), ry + Inches(0.06),
                  col_widths[c_idx] - Inches(0.1), row_h - Inches(0.08),
                  size=12, color=cell_color, align=PP_ALIGN.LEFT)
            cx += col_widths[c_idx]


# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, Inches(2.8), W, Inches(0.08), fill=GOLD)

txbox(sl, "POP Demand & Order Intelligence",
      Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.4),
      size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl, "Hack the Coast 2026  ·  Prince of Peace Enterprises",
      Inches(0.8), Inches(3.05), Inches(11.5), Inches(0.55),
      size=22, color=GOLD, align=PP_ALIGN.CENTER)

txbox(sl, "CPG importer/distributor  ·  Tiger Balm · POP Ginger Chews · Ferrero\n"
          "~800 SKUs  ·  3 DCs (SF / NJ / LA)  ·  No WMS, no API — CSV/Excel only",
      Inches(1.5), Inches(3.85), Inches(10), Inches(0.9),
      size=16, color=LGRAY, align=PP_ALIGN.CENTER)

txbox(sl, "Two deliverables:  Reorder Alert  +  Demand Curve",
      Inches(1.5), Inches(5.0), Inches(10), Inches(0.55),
      size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 2: The Problem ───────────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "The Problem", "How POP buys today — and why it breaks")
footer(sl)

# left column: today's process
rect(sl, Inches(0.4), Inches(1.35), Inches(5.6), Inches(5.6), fill=LGRAY)
txbox(sl, "Today's Manual Process", Inches(0.55), Inches(1.45), Inches(5.3), Inches(0.45),
      size=15, bold=True, color=NAVY)
steps = [
    "Open Excel export of recent sales",
    "Eyeball last few months — guess ~500 cases/mo",
    "Check current inventory",
    "Recall lead time from memory (~90 days China)",
    "Place PO manually — repeat for 1,000 SKUs",
]
bullet_box(sl, steps, Inches(0.55), Inches(1.95), Inches(5.3), Inches(3.5), size=15)

txbox(sl, "No alert system — buyer must remember to check every SKU",
      Inches(0.55), Inches(5.3), Inches(5.3), Inches(0.55),
      size=13, bold=True, color=RED)

# right column: what goes wrong
rect(sl, Inches(6.6), Inches(1.35), Inches(6.3), Inches(5.6), fill=LGRAY)
txbox(sl, "What Goes Wrong", Inches(6.75), Inches(1.45), Inches(6.0), Inches(0.45),
      size=15, bold=True, color=NAVY)
problems = [
    "Promo months included → inflated baseline → over-order",
    "Stockout months included → deflated baseline → under-order",
    "~23% of rows are noise (promo / markdown / stockout)",
    "Safety stock is one-size-fits-all across 800 SKUs",
    "No per-DC visibility — SF, NJ, LA treated as one",
]
bullet_box(sl, problems, Inches(6.75), Inches(1.95), Inches(6.0), Inches(3.5), size=15)

txbox(sl, "Our fix: strip the noise first, then do the math",
      Inches(6.75), Inches(5.3), Inches(6.0), Inches(0.55),
      size=13, bold=True, color=TEAL)


# ── SLIDE 3: Demand Cleaning ──────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "How We Cleaned Demand", "Three filters applied to 236,818 transactions across 3 years")
footer(sl)

col_w = [Inches(2.2), Inches(1.2), Inches(7.4)]
headers = ["Filter", "Share", "Method / Finding"]
rows = [
    ("🔴  Promo / TPR",
     "10.2%",
     "Chargeback file: 6,868 rows → 5,004 real TPRs (73%) after removing publication fees, shelf-talker fees, trade-show charges. Join on Customer × Brand × Month."),
    ("🟠  Markdown",
     "13.3%",
     "Per-(SKU × channel) median price; flag rows < 85% of median. Per-channel threshold crucial — Health Food shelf premiums broke a pooled median."),
    ("🟡  Stockout / Lost Demand",
     "0.04%",
     "Reconstructed 3 yrs of inventory by rewinding from today's snapshot. Nearly all flags = T-32206 SF, 4 weeks in May–Jun 2023. DC substitution masked most stockouts."),
    ("✅  Clean Demand",
     "77.2%",
     "None of the above. This is the organic run-rate signal fed to the reorder math and elasticity curves."),
]

simple_table(sl, headers, rows,
             x=Inches(0.4), y=Inches(1.3),
             col_widths=col_w, row_h=Inches(1.15))

txbox(sl, "Key insight: POP's demand is promo/markdown-polluted (22.8%), not stockout-polluted (0.04%). "
          "The cleaning job is promo separation, not lost-sales imputation.",
      Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.65),
      size=14, bold=True, color=NAVY)


# ── SLIDE 4: Channel Analysis ─────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Channel-Level Demand Analysis", "Same SKU — three very different demand behaviors")
footer(sl)

col_w = [Inches(2.5), Inches(1.3), Inches(1.5), Inches(1.5), Inches(5.9)]
headers = ["Channel", "Rows", "Clean %", "Elasticity β", "What it means"]
rows = [
    ("MM  American Mainstream\n(Target, Walmart, Kroger)",
     "113,226", "67.5%", "−3.33",
     "Most promo-heavy. Retailers run heavy TPR calendars. Fat safety stock needed — volume swings with every promo."),
    ("AM  Asian Ethnic Market",
     "107,399", "85.9%", "−3.96",
     "Cleanest signal. Almost no promos run. Most elastic — price cuts do move volume. Best channel to read organic trend."),
    ("HF  Health Food\n(Whole Foods, Natural Grocers)",
     "16,189", "56.2%", "−1.31",
     "Most polluted AND most inelastic. Distributor-mediated — price changes barely move volume. Don't run TPRs here; you'll just give away margin."),
]

simple_table(sl, headers, rows,
             x=Inches(0.4), y=Inches(1.3),
             col_widths=col_w, row_h=Inches(1.45))

txbox(sl, "Implication: one-size-fits-all safety stock is wrong. Per-channel elasticity drives per-SKU buffer sizing in F1.",
      Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.55),
      size=14, bold=True, color=NAVY)


# ── SLIDE 5: Before / After ───────────────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Before / After: Tiger Balm Patch Warm (T-32206)",
           "Showcase SKU — MM channel, ~17,700 units/week at peak")
footer(sl)

col_w = [Inches(3.5), Inches(3.5), Inches(3.5), Inches(2.4)]
headers = ["Metric", "Raw (before)", "Clean (after)", "Impact"]
rows = [
    ('"Promo" rows',          ("55%", RED),        ("10%", GREEN),   "−45 pp after fee filter"),
    ("Weekly run rate MM-NJ", ("~19,000 u/wk", RED), ("15,348 u/wk", GREEN), "~20% lower — organic only"),
    ("Suggested PO (NJ)",     ("Over-estimated", RED), ("7,571 cases", GREEN), "Right-sized to true demand"),
    ("SF DC alert?",          ("Would fire", RED),  ("Suppressed ✓", GREEN), "SF has 49-wk cover post-2023 dip"),
]

simple_table(sl, headers, rows,
             x=Inches(0.4), y=Inches(1.3),
             col_widths=col_w, row_h=Inches(0.85))

# summary box
rect(sl, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.05), fill=LGRAY)
txbox(sl, "Reorder output across all SKUs:  165 alerts  ·  57 SKUs × 3 DCs  ·  109 high-confidence / 28 medium / 96 low",
      Inches(0.55), Inches(4.98), Inches(12.2), Inches(0.45),
      size=14, bold=True, color=NAVY)
txbox(sl, "Low-confidence = missing lead time or case-pack in item master — flagged for manual review, not suppressed",
      Inches(0.55), Inches(5.43), Inches(12.2), Inches(0.4),
      size=12, color=DGRAY)


# ── SLIDE 6: Backtesting (placeholder) ───────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Reorder Alert Backtesting (F1)", "Validating alert quality against historical outcomes")
footer(sl)

# placeholder banner
rect(sl, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.1), fill=GOLD)
txbox(sl, "⏳  Results pending — placeholder",
      Inches(1.5), Inches(2.25), Inches(10.3), Inches(0.9),
      size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# methodology box
rect(sl, Inches(0.4), Inches(3.65), Inches(12.5), Inches(2.8), fill=LGRAY)
txbox(sl, "Methodology", Inches(0.6), Inches(3.75), Inches(12.0), Inches(0.4),
      size=15, bold=True, color=NAVY)

method_items = [
    "Re-run reorder alerts on historical clean demand (hold out last 6 months)",
    "Compare each fired alert to whether that SKU × DC actually ran low in the subsequent lead-time window",
    "Metric 1 — Precision: of flagged SKUs, what % actually ran below safety stock? (measures false alarms)",
    "Metric 2 — Recall: of real stockout events, what % did we flag in advance? (measures missed alerts)",
    "Metric 3 — Dollar value of avoided stockouts vs. cost of excess safety stock ordered",
]
bullet_box(sl, method_items, Inches(0.6), Inches(4.2), Inches(12.0), Inches(2.1), size=13)


# ── SLIDE 7: Assumptions & Lessons ───────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, W, H, fill=WHITE)
header_bar(sl, "Assumptions, Limitations & Lessons Learned", "")
footer(sl)

# left: limitations
rect(sl, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.55), fill=LGRAY)
txbox(sl, "Known Limitations", Inches(0.55), Inches(1.42), Inches(5.7), Inches(0.4),
      size=15, bold=True, color=RED)
limits = [
    "Promo join granularity is Brand × Month — over-flags adjacent sub-brands in same customer-month (conservative bias, acceptable)",
    "Off-invoice TPR: some retailers use rebate-style TPR (chargeback settles discount, not on-invoice price cut) → elasticity β underestimates true promo response for those accounts",
    "Historical inventory reconstructed, not measured — 2.3% null weeks; all flagged as low-confidence",
    "14 SKUs lack parsed lead time → fall back to 13-week default, flagged low-confidence",
    "No forecasting — organic run rate is a trailing average, not a forward model",
]
bullet_box(sl, limits, Inches(0.55), Inches(1.88), Inches(5.7), Inches(4.8), size=12)

# right: surprises + v2
rect(sl, Inches(6.9), Inches(1.35), Inches(5.9), Inches(2.55), fill=LGRAY)
txbox(sl, "What Surprised Us", Inches(7.05), Inches(1.42), Inches(5.7), Inches(0.4),
      size=15, bold=True, color=TEAL)
surprises = [
    "Fee filtering cut 'promo' rows from 55% → 10% — the single biggest accuracy improvement",
    "Stockouts are nearly irrelevant (0.04%); DC substitution absorbs most supply gaps",
    "Per-channel markdown thresholds — pooled medians were off by 6× in Health Food",
]
bullet_box(sl, surprises, Inches(7.05), Inches(1.88), Inches(5.7), Inches(2.0), size=12)

rect(sl, Inches(6.9), Inches(4.1), Inches(5.9), Inches(2.8), fill=NAVY)
txbox(sl, "V2 Roadmap", Inches(7.05), Inches(4.18), Inches(5.7), Inches(0.4),
      size=15, bold=True, color=GOLD)
v2 = [
    "Web app UI — pipeline is wired, connect a viz layer",
    "SKU-level chargeback join → finer promo tagging",
    "Isotonic regression fallback for kinked demand curves",
    "Stockout lost-demand dollar quantification",
    "Confirm off-invoice TPR treatment with POP president",
]
bullet_box(sl, v2, Inches(7.05), Inches(4.62), Inches(5.7), Inches(2.1), size=12)
# fix v2 bullet color to white
# re-draw over the navy box
for item in v2:
    pass  # handled below via direct paragraph color

# redraw v2 bullets with white text
tb = sl.shapes[-1]
for para in tb.text_frame.paragraphs:
    for run in para.runs:
        run.font.color.rgb = WHITE


out = "/Users/jhuang130/pop_prompt2/presentation/pop_demand_intelligence.pptx"
prs.save(out)
print(f"Saved → {out}")
